"""Generate a stakeholder review PowerPoint presentation for the Inventory Management System."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def set_slide_bg(slide, r, g, b):
    """Set a solid background color on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_title_slide(prs):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, 0x1B, 0x3A, 0x5C)  # Dark blue

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Inventory Management System"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Stakeholder Review"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p2.alignment = PP_ALIGN.CENTER

    # Date
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "August 2026 | Warehouse Operations Engineering"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p3.alignment = PP_ALIGN.CENTER


def add_agenda_slide(prs):
    """Slide 2: Agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xF5, 0xF7, 0xFA)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Agenda items
    items = [
        "System Purpose & Scope",
        "Functional Capabilities",
        "Architecture Overview",
        "Data Model",
        "Current Inventory Snapshot",
        "Reorder Workflow",
        "Reporting Capabilities",
        "Technical Decisions & Trade-offs",
        "Roadmap & Next Steps",
        "Q&A",
    ]

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"{i + 1}.  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)


def add_purpose_slide(prs):
    """Slide 3: System Purpose & Scope."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "System Purpose & Scope"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    bullets = [
        "Lightweight inventory management for a small warehouse operation",
        "Tracks products across categories, manages stock movements, and automates reorder alerts",
        "REST API with JSON responses — integrates with existing tools and UIs",
        "Designed as an internal microservice (no external-facing UI in v1)",
        "In-memory storage for rapid prototyping; persistence planned for v2",
    ]

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(14)


def add_capabilities_slide(prs):
    """Slide 4: Functional Capabilities."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Functional Capabilities"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    capabilities = [
        ("Product Management", "Add, update, remove products (SKU, price, category, thresholds)"),
        ("Category Management", "Organize products; enforced referential integrity on delete"),
        ("Stock Movements", "Record RECEIVED, SHIPPED, ADJUSTMENT, RETURNED with audit trail"),
        ("Reorder Alerts", "Automatic threshold detection with suggested reorder quantities"),
        ("Reporting", "Stock summary, movement history, category-level aggregates"),
    ]

    y = Inches(1.4)
    for title, desc in capabilities:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.4), Inches(0.9))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

        p2 = tf2.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        y += Inches(1.0)


def add_architecture_slide(prs):
    """Slide 5: Architecture Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Architecture Overview"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Draw simplified architecture boxes
    layers = [
        ("Client (HTTP)", 0xE8, 0xF4, 0xFD),
        ("API Layer (Flask)", 0xBE, 0xDA, 0xF7),
        ("Service Layer (Business Logic)", 0x7F, 0xB3, 0xE0),
        ("Model Layer (Pydantic)", 0x4A, 0x90, 0xD9),
        ("Data Store (In-Memory + JSON seed)", 0x1B, 0x3A, 0x5C),
    ]

    y = Inches(1.3)
    for label, r, g, b in layers:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(6), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if b < 0x90 else RGBColor(0x1B, 0x3A, 0x5C)
        p.alignment = PP_ALIGN.CENTER

        y += Inches(1.05)

    # Note
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Each layer communicates only with its immediate neighbor. See architecture.puml for full component diagram."
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x77, 0x77, 0x77)


def add_data_model_slide(prs):
    """Slide 6: Data Model."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Data Model"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Core entities
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Core Entities"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    entities = [
        "Category — id, name, description",
        "Product — SKU, name, category, unit price, reorder threshold",
        "StockMovement — id, product SKU, type, quantity, timestamp, note",
    ]
    for entity in entities:
        p = tf2.add_paragraph()
        p.text = f"\u2022  {entity}"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)

    # Derived
    txBox3 = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(4))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "Computed / Derived"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    derived = [
        "ReorderAlert — generated when stock < threshold",
        "StockSummaryItem — aggregated view with status & value",
        "StockStatus enum — OK, LOW, OUT_OF_STOCK",
        "MovementType enum — RECEIVED, SHIPPED, ADJUSTMENT, RETURNED",
    ]
    for item in derived:
        p = tf3.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)


def add_inventory_snapshot_slide(prs):
    """Slide 7: Current Inventory Snapshot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Current Inventory Snapshot"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Table data
    rows = [
        ("Product", "Category", "Stock", "Threshold", "Status"),
        ("USB-C Cable (2m)", "Electronics", "155", "50", "OK"),
        ("Wireless Mouse", "Electronics", "25", "30", "LOW"),
        ("A4 Paper (500 sheets)", "Office Supplies", "80", "100", "LOW"),
        ("Ballpoint Pen (Black)", "Office Supplies", "500", "200", "OK"),
        ("Shipping Box (Medium)", "Packaging", "120", "150", "LOW"),
        ("Packing Tape Roll", "Packaging", "150", "80", "OK"),
    ]

    table = slide.shapes.add_table(len(rows), 5, Inches(0.3), Inches(1.3), Inches(9.4), Inches(3.5)).table

    # Set column widths
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.2)
    table.columns[3].width = Inches(1.4)
    table.columns[4].width = Inches(1.2)

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)

            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
            else:
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                # Color-code status
                if col_idx == 4 and cell_text == "LOW":
                    p.font.color.rgb = RGBColor(0xE8, 0x6C, 0x00)
                    p.font.bold = True
                elif col_idx == 4 and cell_text == "OK":
                    p.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

    # Summary note
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = "3 of 6 products currently below reorder threshold"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xE8, 0x6C, 0x00)


def add_reorder_workflow_slide(prs):
    """Slide 8: Reorder Workflow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Reorder Workflow"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    steps = [
        "1.  System scans all products (triggered daily or on-demand via API)",
        "2.  Compares current stock level against each product\u2019s reorder threshold",
        "3.  If stock < threshold \u2192 generates a ReorderAlert",
        "4.  Suggested reorder quantity = (threshold \u00d7 2) \u2212 current stock",
        "5.  Alerts available via GET /stock/reorder-alerts",
    ]

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, step in enumerate(steps):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = step
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(12)

    # Example box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.5), Inches(8), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
    shape.line.color.rgb = RGBColor(0x4A, 0x90, 0xD9)
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Example: Wireless Mouse"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    p2 = tf.add_paragraph()
    p2.text = "Current stock: 25 | Threshold: 30 | Suggested order: (30\u00d72) \u2212 25 = 35 units"
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_reporting_slide(prs):
    """Slide 9: Reporting Capabilities."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Reporting Capabilities"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    reports = [
        ("Stock Summary Report", "All products with current level, status (OK / LOW / OUT_OF_STOCK), unit price, and total stock value"),
        ("Movement History", "Full audit trail of all stock movements, filterable by product and movement type"),
        ("Category Summary", "Total products and total stock value aggregated per category"),
    ]

    y = Inches(1.4)
    for title, desc in reports:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.4), Inches(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

        p2 = tf2.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2.space_before = Pt(4)

        y += Inches(1.5)

    # Footer note
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf3 = txBox3.text_frame
    p = tf3.paragraphs[0]
    p.text = "All reports available via REST API endpoints — ready for dashboard integration."
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x77, 0x77, 0x77)


def add_tradeoffs_slide(prs):
    """Slide 10: Technical Decisions & Trade-offs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Technical Decisions & Trade-offs"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    # Table
    rows = [
        ("Decision", "Rationale", "Trade-off"),
        ("In-memory storage", "Simplicity, fast iteration", "No persistence across restarts"),
        ("Flask framework", "Lightweight, well-understood", "Less structure for large projects"),
        ("Pydantic models", "Strong validation, clear contracts", "Additional dependency"),
        ("JSON seed data", "Easy test/demo setup", "Not a production data strategy"),
    ]

    table = slide.shapes.add_table(len(rows), 3, Inches(0.3), Inches(1.3), Inches(9.4), Inches(3)).table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(3.4)

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)

            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
            else:
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_roadmap_slide(prs):
    """Slide 11: Roadmap & Next Steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Roadmap & Next Steps"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    sections = [
        ("Short-term", [
            "Add persistent storage (SQLite or PostgreSQL)",
            "Implement authentication & role-based access",
            "Date-range filtering on movement history",
        ]),
        ("Medium-term", [
            "Dashboard UI for warehouse staff",
            "Automated reorder email/webhook notifications",
            "Batch import/export (CSV)",
        ]),
        ("Long-term", [
            "Multi-warehouse support",
            "Barcode/QR scanning integration",
            "Demand forecasting with historical data",
        ]),
    ]

    y = Inches(1.2)
    for heading, items in sections:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.4), Inches(1.6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        p.text = heading
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x4A, 0x90, 0xD9)

        for item in items:
            p = tf2.add_paragraph()
            p.text = f"\u2022  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(4)

        y += Inches(1.75)


def add_qa_slide(prs):
    """Slide 12: Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0x1B, 0x3A, 0x5C)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Thank you for your time"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p2.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_agenda_slide(prs)
    add_purpose_slide(prs)
    add_capabilities_slide(prs)
    add_architecture_slide(prs)
    add_data_model_slide(prs)
    add_inventory_snapshot_slide(prs)
    add_reorder_workflow_slide(prs)
    add_reporting_slide(prs)
    add_tradeoffs_slide(prs)
    add_roadmap_slide(prs)
    add_qa_slide(prs)

    output_path = "Inventory_Management_Stakeholder_Review.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
